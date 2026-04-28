from flask import Flask, request, render_template, send_file, redirect, url_for, session, after_this_request
import secrets
import os
import re
import pytesseract
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from openai import OpenAI
from PIL import Image
from pydantic import BaseModel
import json
import ast
import time
from concurrent.futures import ThreadPoolExecutor, as_completed
import unicodedata
import uuid
import sys
import sys, os
import tempfile
import shutil
import os, sys
import io, sys

import sys, os

log_path = os.path.join(os.path.expanduser("~"), "ScriptFormatter_log.txt")
try:
    log_file = open(log_path, "a", encoding="utf-8", buffering=1)
except Exception:
    class _NullWriter:
        def write(self, *_): pass
        def flush(self): pass
    log_file = _NullWriter()

sys.stdout = log_file
sys.stderr = log_file

print(f"[DEBUG] Logging started. Writing to: {log_path}")


APP_NAME = "ScriptFormatter"

if hasattr(sys, "_MEIPASS"):
    BASE_DIR = sys._MEIPASS  # templates & static live here (read-only)
else:
    BASE_DIR = os.path.abspath(".")

# Use user's local app data for working directories
if sys.platform == "win32":
    WORK_DIR = os.path.join(os.getenv("LOCALAPPDATA", os.path.expanduser("~")), APP_NAME)
else:
    WORK_DIR = os.path.join(os.path.expanduser("~/.local/share"), APP_NAME)

# Make sure folders exist
for sub in ["uploads", "presentations", "logs", "output_paths"]:
    os.makedirs(os.path.join(WORK_DIR, sub), exist_ok=True)

print(f"[DEBUG] BASE_DIR = {BASE_DIR}")
print(f"[DEBUG] WORK_DIR = {WORK_DIR}")



app = Flask(
    __name__,
    template_folder=os.path.join(BASE_DIR, "templates"),
    static_folder=os.path.join(BASE_DIR, "static")
)

app.secret_key = os.environ.get("FLASK_SECRET_KEY", "dev-default-key")

if sys.platform == "win32":
    TESSERACT_PATH = os.path.join(BASE_DIR, "bin", "tesseract", "tesseract.exe")
    POPPLER_PATH = os.path.join(BASE_DIR, "bin", "poppler", "Library", "bin")
elif sys.platform == "darwin":
    TESSERACT_PATH = os.path.join(BASE_DIR, "bin", "tesseract", "tesseract")  # mac binary, no .exe
    POPPLER_PATH = os.path.join(BASE_DIR, "bin", "poppler")                   # poppler binaries live here
else:
    # fallback for Linux (optional)
    TESSERACT_PATH = "tesseract"
    POPPLER_PATH = None

pytesseract.pytesseract.tesseract_cmd = TESSERACT_PATH

def pdf_to_images(pdf_path, dpi=300):
    kwargs = {"dpi": dpi, "fmt": "png"}
    if POPPLER_PATH:
        kwargs["poppler_path"] = POPPLER_PATH
    return convert_from_path(pdf_path, **kwargs)


def split_text(text, sentences_per_chunk=50):
    sentences = re.split(r'(?<=[.!?])\s+', text.strip())
    chunks = [
        ' '.join(sentences[i:i + sentences_per_chunk])
        for i in range(0, len(sentences), sentences_per_chunk)
    ]
    return chunks

def extract_text_from_images(images):
    extracted_text = []
    for img in images:
        img = preprocess_image(img)
        text = pytesseract.image_to_string(img, config="--psm 6")
        extracted_text.append(text.strip())

    raw_text = "\n".join(extracted_text)

    return raw_text

def preprocess_image(image):
    return image.convert("L")

def recover_valid_json_entries(raw_content: str):
    entries = []
    # Try to extract inside the "entries": [ ... ]
    match = re.search(r'"entries"\s*:\s*\[(.*)', raw_content, re.DOTALL)
    if not match:
        print("No entries array found in raw content.")
        return []

    entries_raw = match.group(1)

    # Extract possible objects using a greedy balanced-brace parser
    brace_stack = 0
    current_obj = ""
    for c in entries_raw:
        current_obj += c
        if c == '{':
            brace_stack += 1
        elif c == '}':
            brace_stack -= 1

        if brace_stack == 0 and current_obj.strip():
            try:
                entry = json.loads(current_obj)
                entries.append(entry)
            except Exception:
                # Fallback: try to salvage fields manually
                salvaged = recover_fields_from_partial(current_obj)
                if salvaged:
                    entries.append(salvaged)
            current_obj = ""

    # Also handle trailing unbalanced object (e.g., truncated at end)
    if brace_stack > 0 and current_obj.strip():
        salvaged = recover_fields_from_partial(current_obj)
        if salvaged:
            entries.extend(salvaged)

    return entries

def recover_fields_from_partial(fragment: str):
    try:
        type_match = re.search(r'"type"\s*:\s*"(\w+)"', fragment)
        character_match = re.search(r'"character"\s*:\s*"([^"]*)"', fragment)
        text_match = re.search(r'"text"\s*:\s*"((?:[^"\\]|\\.)*)', fragment)  # captures even if no closing quote

        if type_match and text_match:
            entry = {
                "type": type_match.group(1),
                "character": character_match.group(1) if character_match else "",
                "text": text_match.group(1).encode('utf-8').decode('unicode_escape')
            }
            return [
                entry
            ]
    except Exception as e:
        print("Could not recover partial object:", e)
    return None

def format_section(index, section, client, prompt, schema):
    try:
        response = client.chat.completions.create(
            model="gpt-4.1-mini",
            messages=[
                {"role": "system", "content": prompt},
                {"role": "user", "content": section}
            ],
            response_format={"type": "json_schema", "json_schema": schema}
        )
        content = response.choices[0].message.content
        try:
            parsed = json.loads(content)
            listed_json = parsed.get("entries", [])
            listed_json.append({
                    "type": "system",
                    "character": "SYSTEM NOTE",
                    "text": f"Possible line error. Check for missing lines."
                })
            return (index, listed_json)
        except json.JSONDecodeError:
            recovered = recover_valid_json_entries(content)
            if recovered:
                recovered.append({
                    "type": "system",
                    "character": "SYSTEM NOTE",
                    "text": f"BLOCK {index} had partial error due to {response.choices[0].finish_reason}, recovered entries, but likely missing lines. Manually fill in"
                })
                return (index, recovered)
            else:
                return (index, [{
                    "type": "system",
                    "character": "SYSTEM NOTE",
                    "text": f"BLOCK {index} unrecoverable, missing 500 words. Manually fill in."
                }])
    except Exception as e:
        return (index, [{
            "type": "system",
            "character": "SYSTEM NOTE",
            "text": f"BLOCK {index} failed: {str(e)}"
        }])
    
def clean_and_structure_script(text, api_key, upload_id):
    client = OpenAI(api_key=api_key)
    text_sections = split_text(text, sentences_per_chunk=50)  # e.g., into 500-word chunks
    prompt = """
        You are a professional script formatter. Your task is to convert raw, unformatted theatrical script text into a clean structured format for a slide generator.

        Follow these rules exactly and strictly:
        
        1. Censorship and content safety:
        - IMPORTANT: If the script contains any of the following types of sensitive or restricted content, censor only the individual offending words using asterisks (e.g., “****”). Do not remove, paraphrase, skip, or alter the rest of the line. Maintain line structure, speaker names, and meaning.
        - You must censor:
            -Profanity and vulgar language (e.g., f***, s***, b****)
            -Slurs or hate speech (racial, ethnic, gender, religious, homophobic, etc.)
            -Sexual content:
            -Explicit sexual terms or language
            -Sexual innuendo or suggestive descriptions
            -Sexual euphemisms, slang, or references to body parts
            -Violence and gore:
            -Graphic descriptions of injury or harm
            -Excessive violent language (e.g., “beat her to a pulp”)
            -Self-harm or suicide references
            -Drug and alcohol abuse:
            -Mentions of illegal drugs, addiction, or overdosing
            -Abusive or excessive alcohol use
            -Derogatory or dehumanizing language:
            -Terms meant to belittle, demean, or objectify individuals or groups
            -Also censor borderline content even if not explicit, including:
            -Common euphemisms for explicit content (e.g., “doing it”, “getting laid”)
            -Suggestive language in song lyrics if sexual in nature
            -Slang with a double meaning, if clearly used in a sexual or offensive context
            -If you are unsure whether a term is offensive or filtered: err on the side of censoring the word with asterisks.
        - Do not skip or truncate these lines — always keep the dialogue intact, except for censoring specific words.

        2.  If unsure whether to output a line, apply this rule:
        - If the line could plausibly be spoken or sung aloud by a character, output it.
        - If it clearly describes an action, is contained by parenthesis, or is a cue or direction, omit it.
        - When uncertain between dialogue and stage direction, only output if the line includes a full sentence, emotion, or intent that a character might express aloud.
        - Do not output parentheticals, even if they contain tone or direction.

        3. ALWAYS omit all stage directions, including:  
        - Lines describing actions (e.g., "SKY enters", "She crosses the stage", "beat", etc.)  
        - Parenthetical text (e.g., "(quietly)", "(aside)", "(reading)")  
        - Sound cues (e.g., "MUSIC PLAYS", "Blackout", "Lights fade")
        - Check that your output contains none of these cases
        - Blocks of text may often start with stage directions. ALWAYS omit these.
        - Be especially cautious of lines starting or ending with parentheses, stage actions (like "She walks over"), or technical cues. These are almost never dialogue or singing.

        4. Scene titles:  
        OMIT all scene titles; do not output any entry of type 'scene'

        5. Dialogue:  
        Lines with the format "CHARACTER: Line of text", or a name of a CHARACTER, followed by a length of text, are treated as spoken dialogue:  
        { "type": "dialogue", "character": "CHARACTER", "text": "line here" } 
        - There may often be stage directions contained within a character's line. In this case, keep the line, but omit the stage directions
        - If a line starts with a capitalized name followed by a sentence—even without a colon—treat it as a dialogue line, as long as it's not stage directions or a song lyric. 
        - Always begin a new entry when a colon is used after a name.  
        - When not sure if text is dialouge or not, output as dialogue. Only omit if POSITIVE it is not meant to be spoken on stage
        - Join multi-character names with slashes: "LISA / ALI".  
        - Correct all spelling and grammar errors in both the character names and the spoken dialogue.
        - If a dialogue line contains parentheticals (e.g., “(angrily)” or “(beat)”), omit those parts but keep the rest of the spoken line.
        - If a line looks like a stage direction without a clear speaker or quote, omit it.

        6. Singing (ALL CAPS lines):
        - Any line that is in ALL CAPS and is not a stage direction must be treated as singing. Use the following rules to structure and assign these lines:

        A. Detecting Singing:
        - If a line is in ALL CAPS and not clearly a cue, direction, or title, always treat it as singing.
        - Skip ALL-CAPS lines that are clearly stage cues, scene headings, or sound/light directions (e.g., "BLACKOUT", "SCENE 2", "MUSIC PLAYS").
        - Do NOT include standalone song titles (e.g., "HONEY HONEY", "DANCING QUEEN") — skip them.

        B. Segmenting Long Singing Passages:
        - If a character is singing a long line with no punctuation, split it into multiple entries at natural phrase boundaries (e.g., a pause, breath, or change in idea).
        - Each phrase must be output as a separate entry.
        - Use the same "character" as the entry it was split from and set the "type" to "singing" for each part.

        C. Character Detection for Singing:
        - ALL-CAPS singing lines often do not follow the "CHARACTER: LINE" format. Some may just begin with a name (e.g., "ALICE").
        - If the line starts with a single word or a short list of 1-3 slash-separated names (e.g., "ALICE", "LISA / ALI"), treat it as a character or group name and assign it as the "character" for that line.
        - Join multi-character names with slashes, like: "LISA / ALI". 
        - Character names should never exceed 3 words.
        - If an ALL-CAPS line has 4 or more words and does not look like a character name, assume it is pure lyrics. Assign it to the most recent valid singing character.
        - If a new character/group name appears mid-line, it signals a new singer. Begin a new "singing" entry with that character.

        D. Character Validity Check:
        - Before outputting any singing entry, verify the "character" name:
        - Only use it if it is an actual character or group name.
        - If it is not valid or unknown, default to the most recent singing character.
        - Do not invent or guess lyrics for singing lines. Use only what's in the source.

        E. Output Format:
        - All "text" values for singing must remain in ALL CAPS.
        - Valid singing entries must follow this format:
        { "type": "singing", "character": "CHARACTER", "text": "SUNG LINE" }

        7. Continuations:  
        If the script begins midway through a character's line (i.e., no character name at the start of the first line), but it is NOT stage directions, return that line as:  
        { "type": "continuation", "text": "remaining text here" }  
        - Continuations should ALWAYS only be the first JSON entry of your output. Every other entry should be dialogue or singing. If you believe something is a continuation, but it is not the first entry, it is likely a stage direction; omit it.
        - Do NOT assign a character name.  
        - Once a new Character or Group is singing/speaking, ALWAYS create a new entry with the correct formatting.
        - Continue until a new Character or Group name (ie ENSEMBLE, COMPANY, or CAST), at which point check if you should begin a new singing or dialogue line.
        - If the first line of the input is in ALL CAPS and does not look like a sound cue or heading, treat it as a sung lyric. Use the next nearby character or group name to assign the singer retroactively if needed.
        - Continuation lines should not be super long, if they are it indicates you forgot to begin a new entry for a singing or dialogue line. Check, and start new entries where you were supposed to.


        8. Process line by line, strictly in order. Do not merge lines, infer missing names, or reorder content.

        9. Correct all typos and grammatical errors in the input script, including dialogue text, character names, and sung lines. The output must be fully clean and readable.

        9a. Additional typo-fix rules:
        - Detect and split run-together repeats: if a word or lyric is repeated without spaces, restore the spaces.  
            e.g. “I DOIDOIDO, I DOIDO” → "I DO I DO I DO, I DO I DO"
        - Correct other fused words or doubled letters (e.g. "HELLOWORLD" → "HELLO WORLD").

        10. End of show:
        - Even if the section indicates the show is over, there may be an encore. Continue to format until the end of every section.

        11. Only return valid structured entries in this JSON format:
        {
        "entries": [
            { "type": "continuation", "text": "..." },
            { "type": "singing", "character": "...", "text": "..." },
            { "type": "dialogue", "character": "...", "text": "..." }
        ]
        }

        Example input:
            and I never saw her agan.  
            MUSIC PLAYS.  
            SOPIHE: But you loved her.  
            "HONEY HONEY"  
            LISA AND ALI 
            JUST SAY IT!
            WOMEN
            BECAUSE IT'S TRUE IDO IDOI DOIDOIDO

        Example output:
        {
        "entries": [
            { "type": "continuation", "text": "and I never saw her again." },
            { "type": "dialogue", "character": "SOPHIE", "text": "But you loved her." },
            { "type": "singing", "character": "LISA / ALI", "text": "JUST SAY IT!" },
            { "type": "singing", "character": "WOMEN", "text": "BECAUSE IT'S TRUE" },
            { "type": "singing", "character": "WOMEN", "text": "I DO I DO I DO I DO I DO" }
        ]
        }

        IMPORTANT: Return ONLY a single valid JSON object matching the schema. No explanations or extra text.
        """
    schema = {
            "name": "script_entries",
            "strict": True,
            "schema": {
                "type": "object",
                "properties": {
                "entries": {
                    "type": "array",
                    "items": {
                    "type": "object",
                    "properties": {
                        "type": {
                        "type": "string",
                        "enum": ["dialogue", "continuation", "singing"]
                        },
                        "character": {
                        "type": ["string", "null"]
                        },
                        "text": {
                        "type": ["string", "null"]
                        }
                    },
                    "required": ["type", "character", "text"],
                    "additionalProperties": False
                    }
                }
                },
                "required": ["entries"],
                "additionalProperties": False
            }
            }

    results = []
    with ThreadPoolExecutor(max_workers=20) as executor:
        futures = [
            executor.submit(format_section, i, section, client, prompt, schema)
            for i, section in enumerate(text_sections)
        ]
        for future in as_completed(futures):
            results.append(future.result())

    # Reorder based on original index
    results.sort(key=lambda x: x[0])
    ordered_entries = []
    for index, entries in results:
        ordered_entries.extend(entries)

    # Postprocess singing lines, etc.
    final_json = []
    for entry in ordered_entries:
        if entry["type"] == "singing" and entry["text"]:
            split_lines = split_singing(entry["text"], max_words=100)
            for line in split_lines:
                final_json.append({
                    "type": "singing",
                    "character": entry["character"],
                    "text": line
                })
        else:
            final_json.append(entry)
    for entry in final_json:
        if entry.get("text"):
            entry["text"] = normalize_quotes(entry["text"])
        if entry.get("character"):
            entry["character"] = normalize_quotes(entry["character"])
    print(final_json)
    output_dir = os.path.join(WORK_DIR, "logs")
    os.makedirs(output_dir, exist_ok=True)
    json_path = os.path.join(output_dir, f"{upload_id}.json")

    with open(json_path, "w", encoding="utf-8") as f:
        final_json = json.loads(json.dumps(final_json, ensure_ascii=False))
        json.dump(final_json, f, indent=2, ensure_ascii=False)
    return final_json

def fix_mojibake(s: str) -> str:
    try:
        # take the bytes that were wrongly decoded as Latin-1 and
        # decode them correctly as UTF-8
        return s.encode('latin1').decode('utf8')
    except Exception:
        return s
    
def normalize_quotes(s: str) -> str:
    # first undo common mojibake
    s = fix_mojibake(s)

    # then fold all curly/fancy quotes to plain ASCII
    s = unicodedata.normalize('NFKC', s)

    replacements = {
        '“': '"', '”': '"',
        '‘': "'", '’': "'",
        '‚': "'", '‛': "'",
        '`':  "'", '´':  "'",
        # …and if any stray mojibake remains, catch it here:
        'â': "'",  # literal three-char sequence
        'â': '"',
        'â': '"',
    }
    for bad, good in replacements.items():
        s = s.replace(bad, good)
    return s

def split_dialogue(dialogue, max_sentences=5):
    sentences = re.split(r'(?<=[.!?]) +', dialogue)
    return [" ".join(sentences[i:i + max_sentences]) for i in range(0, len(sentences), max_sentences)]

def split_singing(text, max_words=100):
    words = text.split()
    chunks = []

    for i in range(0, len(words), max_words):
        chunk = " ".join(words[i:i + max_words])
        chunks.append(chunk)

    return chunks

def create_presentation(script_data, output_pptx, upload_id):
    prs = Presentation()
    buffered_entry = None         # Holds current dialogue/singing block
    previous_line = None          # Stores (character, full_text) for gray area

    def flush_buffered_entry():
        nonlocal buffered_entry, previous_line
        if not buffered_entry:
            return

        character = buffered_entry.get("character", "")
        full_text = buffered_entry["text"]

        for chunk in split_dialogue(full_text, max_sentences=5):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)

            # Show previous line
            if previous_line:
                prev_char, prev_text = previous_line
                prev_textbox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(2))
                prev_frame = prev_textbox.text_frame
                prev_frame.text = f"{prev_char}: {prev_text}" if prev_char else prev_text
                prev_frame.paragraphs[0].font.color.rgb = RGBColor(89, 89, 89) # Changed
                prev_frame.paragraphs[0].font.size = Pt(25) # Changed
                prev_frame.paragraphs[0].font.name = 'Calibri'
                prev_frame.paragraphs[0].font.bold = True
                prev_frame.word_wrap = True

            # Current line (white)
            cur_textbox = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(2))
            cur_frame = cur_textbox.text_frame
            cur_frame.text = f"{character}: {chunk}" if character else chunk
            cur_frame.paragraphs[0].font.size = Pt(30) # Changed
            cur_frame.paragraphs[0].font.name = 'Calibri' # Changed
            cur_frame.paragraphs[0].font.bold = True
            cur_frame.paragraphs[0].font.color.rgb = RGBColor(166, 166, 166) # Changed
            cur_frame.word_wrap = True

            # Update previous_line
            previous_line = (character, chunk)

        buffered_entry = None  # Clear after flushing

    for entry in script_data:
        entry_type = entry["type"]
        text = entry["text"]

        # Handle continuation
        if entry_type == "continuation" and buffered_entry:
            buffered_entry["text"] += " " + text
            continue

        # If we reach a new top-level entry, flush the buffer first
        if entry_type in ["dialogue", "singing"]:
            flush_buffered_entry()
            buffered_entry = {
                "type": entry_type,
                "character": entry.get("character", ""),
                "text": text
            }

        elif entry_type == "system":
            character = entry.get("character", "")
            text = entry["text"]
            for chunk in split_dialogue(text):
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)
                if previous_line:
                    prev_textbox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(2))
                    prev_frame = prev_textbox.text_frame
                    prev_frame.text = f"{previous_line[0]}: {previous_line[1]}"
                    prev_frame.paragraphs[0].font.color.rgb = RGBColor(89, 89, 89)
                    prev_frame.paragraphs[0].font.name = 'Calibri'
                    prev_frame.paragraphs[0].font.size = Pt(25)
                    prev_frame.paragraphs[0].font.bold = True
                    prev_frame.word_wrap = True
                cur_textbox = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(2))
                cur_frame = cur_textbox.text_frame
                if character:
                    cur_frame.text = f"{character}: {chunk}"
                    previous_line = (character, chunk)
                else:
                    cur_frame.text = chunk
                    previous_line = ("", chunk)
                cur_frame.paragraphs[0].font.size = Pt(30)
                cur_frame.paragraphs[0].font.name = 'Calibri'
                cur_frame.paragraphs[0].font.bold = True
                cur_frame.paragraphs[0].font.color.rgb = RGBColor(166, 166, 166)
                cur_frame.word_wrap = True

        else:
            continue  # skip unsupported types

    # Final flush at end
    flush_buffered_entry()

    prs.save(output_pptx)
    with open(os.path.join(WORK_DIR, "output_paths", f"{upload_id}.txt"), 'w', encoding="utf-8") as file:
        file.write(output_pptx)



@app.route('/', methods=['GET', 'POST'])
def upload_file():
    import traceback
    try:
        if request.method == 'POST':
            upload_id = str(uuid.uuid4())
            pdf_file = request.files['pdf_file']
            api_key = request.form['api_key']
            if pdf_file and api_key:
                pdf_path = os.path.join(WORK_DIR, "uploads", f"{upload_id}.pdf")
                output_pptx = os.path.join(WORK_DIR, "presentations", f"{upload_id}.pptx")
                os.makedirs(os.path.dirname(output_pptx), exist_ok=True)
                output_path_record = os.path.join(WORK_DIR, "output_paths", f"{upload_id}.txt")
                with open(output_path_record, "w", encoding="utf-8") as f:
                    f.write(output_pptx)

                pdf_file.save(pdf_path)
                try:
                    images = pdf_to_images(pdf_path)
                except Exception as e:
                    import traceback
                    with open("error_log.txt", "w", encoding="utf-8") as f:
                        f.write(traceback.format_exc())
                    return f"Error in pdf_to_images: {str(e)}", 500
                text = extract_text_from_images(images)
                script_json = clean_and_structure_script(text, api_key, upload_id)
                if not script_json:
                    return "Error: Could not extract dialogue data", 400
                create_presentation(script_json, output_pptx, upload_id)
                return redirect(url_for('download_page', upload_id=upload_id))
        return render_template('upload.html')
    except Exception as e:
        with open("error_log.txt", "w", encoding="utf-8") as f:
            f.write(traceback.format_exc())
        return "Internal error: see error_log.txt", 500

@app.route('/download_page')
def download_page():
    upload_id = request.args.get('upload_id')
    return render_template('download.html', upload_id=upload_id)

from flask import after_this_request, send_file

@app.route('/download/<upload_id>')
def download_file(upload_id):
    output_pptx = os.path.join(WORK_DIR, "presentations", f"{upload_id}.pptx")

    if not os.path.exists(output_pptx):
        print(f"[ERROR] File not found: {output_pptx}")
        return "File not found", 404

    @after_this_request
    def cleanup(response):
        try:
            os.remove(output_pptx)
            os.remove(os.path.join(WORK_DIR, "uploads", f"{upload_id}.pdf"))
            os.remove(os.path.join(WORK_DIR, "output_paths", f"{upload_id}.txt"))
        except Exception as e:
            print(f"[CLEANUP ERROR] {e}")
        return response

    return send_file(output_pptx, as_attachment=True, download_name=f"{upload_id}.pptx")




if __name__ == "__main__":
    app.run(debug=True)